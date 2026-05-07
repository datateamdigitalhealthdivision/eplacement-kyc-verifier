"""Build editable PowerPoint versions of manuscript Figures 1-4."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
OUTPUT_PPTX = HERE / "editable_figures_1_4.pptx"

COLORS = {
    "ink": RGBColor(23, 32, 51),
    "muted": RGBColor(91, 102, 122),
    "line": RGBColor(190, 202, 222),
    "pale": RGBColor(244, 248, 252),
    "purple": RGBColor(68, 1, 84),
    "blue": RGBColor(49, 104, 142),
    "teal": RGBColor(33, 145, 140),
    "green": RGBColor(53, 183, 121),
    "lime": RGBColor(93, 200, 99),
    "yellow": RGBColor(253, 231, 37),
    "white": RGBColor(255, 255, 255),
}


def set_run(run, size: int, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)
    return box


def title(slide, text, subtitle):
    add_textbox(slide, 0.45, 0.28, 12.4, 0.38, text, size=26, bold=True)
    add_textbox(slide, 0.47, 0.74, 12.2, 0.3, subtitle, size=13, color=COLORS["muted"])


def card(slide, x, y, w, h, heading, body, accent, fill=COLORS["pale"], head_size=16, body_size=12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1.2)
    accent_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.color.rgb = accent
    add_textbox(slide, x + 0.22, y + 0.22, w - 0.35, 0.32, heading, size=head_size, bold=True)
    add_textbox(slide, x + 0.22, y + 0.66, w - 0.35, h - 0.75, body, size=body_size, color=COLORS["muted"])
    return shape


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = COLORS["line"]
    conn.line.width = Pt(2.0)
    conn.line.end_arrowhead = True
    return conn


def pill(slide, x, y, w, h, text, fill=RGBColor(224, 242, 254), color=RGBColor(7, 89, 133)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(186, 230, 253)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, 11, True, color)
    return shape


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "Administrative evidence verification problem", "People applying for placement submit structured claims and heterogeneous proof bundles.")
    items = [
        ("Applicant claims", "Structured form data records marital, health, family, location, disability and examination claims.", COLORS["blue"]),
        ("Document bundles", "Uploaded PDFs mix certificates, letters, medical notes, forms, scans and unrelated pages.", COLORS["lime"]),
        ("Verifier method", "A local-first harness checks only claimed categories using rules, OCR, page images and a VLM.", COLORS["teal"]),
        ("Operator queue", "Reviewers receive proof summaries, page references, confidence and a check/no-check decision.", COLORS["green"]),
    ]
    x = 0.45
    for i, (head, body, accent) in enumerate(items):
        card(slide, x, 1.55, 2.75, 1.95, head, body, accent, head_size=15, body_size=10)
        if i < len(items) - 1:
            arrow(slide, x + 2.85, 2.52, x + 3.2, 2.52)
        x += 3.15
    metrics = [("238", "applicants"), ("1,428", "evidence decisions"), ("6", "proof domains"), ("26.9%", "manual-review queue")]
    for i, (value, label) in enumerate(metrics):
        cx = 1.55 + i * 3.15
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.47), Inches(4.35), Inches(0.95), Inches(0.95))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(236, 253, 245)
        circle.line.color.rgb = RGBColor(153, 246, 228)
        add_textbox(slide, cx - 0.58, 4.55, 1.16, 0.23, value, size=18, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)
        add_textbox(slide, cx - 0.85, 4.86, 1.7, 0.3, label, size=9, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, 6.0, 11.2, 0.45, "Contribution: a claim-guided proof verifier that treats accuracy, auditability, privacy and review workload as one implementation problem.", size=15, bold=True, align=PP_ALIGN.CENTER)


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "Six proof categories used by the verifier", "Evidence is checked only when the applicant claimed that category.")
    items = [
        ("Marriage", "Marriage or nikah certificate; clear spouse relationship proof.", COLORS["blue"]),
        ("Self illness", "Medical proof where the patient is the applicant or strongly implied to be the applicant.", COLORS["teal"]),
        ("Family illness", "Medical proof for spouse, child, parent, dependent or other relevant family member.", COLORS["lime"]),
        ("Spouse location", "Spouse workplace, posting, residence or location evidence relevant to placement.", COLORS["purple"]),
        ("OKU self/family", "Official OKU card, JKM document, disability registration or clear disability evidence.", COLORS["green"]),
        ("MedEX / exam", "MedEX, postgraduate or specialist exam evidence; not routine physical examination.", COLORS["yellow"]),
    ]
    for idx, (head, body, accent) in enumerate(items):
        row = idx // 3
        col = idx % 3
        card(slide, 0.55 + col * 4.25, 1.42 + row * 2.15, 3.65, 1.55, head, body, accent, head_size=15, body_size=10)
    pill(slide, 1.65, 6.25, 10.1, 0.45, "Design rule: unclaimed categories are not labelled as positive evidence.")


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "Claim-guided document evidence workflow", "Operational pipeline from applicant form to auditable decision outputs.")
    steps = [
        ("1. Ingest", "Map applicant columns"),
        ("2. Claims", "Create claimed_* booleans"),
        ("3. PDF", "Acquire original upload"),
        ("4. OCR", "Direct text, render, OCR"),
        ("5. Signals", "Keywords, names, layout"),
        ("6. AI harness", "Verify claimed proof only"),
        ("7. Review", "Structured audit output"),
    ]
    x = 0.35
    for i, (head, body) in enumerate(steps):
        card(slide, x, 1.55, 1.58, 1.65, head, body, [COLORS["blue"], COLORS["teal"], COLORS["lime"]][i % 3], head_size=12, body_size=9)
        if i < len(steps) - 1:
            arrow(slide, x + 1.63, 2.37, x + 1.82, 2.37)
        x += 1.85
    add_textbox(slide, 0.9, 4.25, 11.6, 0.35, "Safeguards embedded across the workflow", size=17, bold=True, align=PP_ALIGN.CENTER)
    safeguards = [
        "Skip unclaimed categories",
        "Cache OCR/model outputs",
        "Second-pass verification",
        "Structured JSON parsing",
        "Proof-strength scoring",
        "Original links retained",
        "Manual review triggers",
    ]
    for i, item in enumerate(safeguards):
        pill(slide, 1.0 + (i % 4) * 2.85, 5.0 + (i // 4) * 0.72, 2.35, 0.42, item)


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "AI harness and local inference architecture", "Model calls are constrained by deterministic context, schema and audit policy.")
    layers = [
        ("Langflow orchestration layer", "Readable workflow graph coordinates loader, fetcher, OCR router, verifier and export writer.", COLORS["blue"]),
        ("Deterministic guardrails", "Claim extraction, category skipping, keyword priors, proof rules and confidence thresholds.", COLORS["teal"]),
        ("Document perception layer", "PDF rendering, direct text extraction, OCR fallback and page-image retention.", COLORS["lime"]),
        ("Local multimodal model layer", "Ollama serves Qwen2.5-VL locally for visual verification of claimed evidence.", COLORS["purple"]),
        ("Structured adjudication layer", "Schema parsing, proof strength, supporting page, summary and review policy.", COLORS["green"]),
    ]
    y = 1.25
    for i, (head, body, accent) in enumerate(layers):
        card(slide, 0.65, y, 8.4, 0.85, head, body, accent, head_size=13, body_size=9)
        if i < len(layers) - 1:
            arrow(slide, 4.85, y + 0.89, 4.85, y + 1.12)
        y += 1.1
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.55), Inches(1.35), Inches(3.2), Inches(4.9))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(236, 253, 245)
    panel.line.color.rgb = RGBColor(153, 246, 228)
    add_textbox(slide, 9.85, 1.75, 2.7, 0.35, "Why this matters", size=17, bold=True, color=COLORS["teal"])
    notes = [
        "No external document API required",
        "Prompt sees applicant claims",
        "Model cannot reward unclaimed categories",
        "Outputs are auditable, not just labels",
        "Operator can inspect source PDF",
    ]
    for i, note in enumerate(notes):
        add_textbox(slide, 9.85, 2.35 + i * 0.65, 2.45, 0.42, f"- {note}", size=10, color=COLORS["ink"])


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    prs.save(OUTPUT_PPTX)
    print(f"Wrote {OUTPUT_PPTX}")


if __name__ == "__main__":
    build()
